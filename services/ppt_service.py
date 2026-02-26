import copy
import hashlib
import json
import uuid
from pathlib import Path

from pptx import Presentation

from utils.logger import logger

TEMPLATE_DIR = Path("storage/templates")
SKELETONS_DIR = Path("storage/layouts/skeletons")
SKELETONS_DIR.mkdir(parents=True, exist_ok=True)


def _copy_slide(source_slide, target_prs: Presentation) -> None:
    """Copy a slide (XML content + relationships) from source into target."""
    slide_layout = source_slide.slide_layout

    target_layout = target_prs.slide_layouts[0]
    for layout in target_prs.slide_layouts:
        if layout.name == slide_layout.name:
            target_layout = layout
            break

    new_slide = target_prs.slides.add_slide(target_layout)

    source_tree = source_slide.shapes._spTree  # noqa: SLF001
    target_tree = new_slide.shapes._spTree     # noqa: SLF001

    for child in list(target_tree):
        target_tree.remove(child)

    for child in source_tree:
        target_tree.append(copy.deepcopy(child))

    for rel in source_slide.part.rels.values():
        if "slideLayout" in rel.reltype:
            continue
        try:
            new_slide.part.rels[rel.rId]
        except KeyError:
            new_slide.part.add_relationship(rel.reltype, rel._target)  # noqa: SLF001


def generate_ppt(skeleton_hash: str, slide_names: list[str]) -> str:
    """
    Merge individual template PPTX files into one skeleton presentation.

    Args:
        skeleton_hash: Deterministic hash used as the output filename.
        slide_names:   Ordered list of template names (without .pptx extension).
                       May include subdirectory prefix, e.g. "PPG/Front".

    Returns:
        Path to the generated skeleton PPTX file.
    """
    if not slide_names:
        raise ValueError("slide_names must not be empty.")

    template_paths: list[Path] = []
    for name in slide_names:
        path = TEMPLATE_DIR / f"{name}.pptx"
        if not path.exists():
            raise FileNotFoundError(f"Template not found: {path}")
        template_paths.append(path)

    logger.info(f"Merging {len(template_paths)} templates into skeleton '{skeleton_hash[:10]}'.")

    final_path = SKELETONS_DIR / f"{skeleton_hash}.pptx"
    if final_path.exists():
        logger.info(f"Skeleton already on disk: {final_path}")
        return str(final_path)

    # Bootstrap from first template to inherit dimensions + theme.
    final_prs = Presentation(str(template_paths[0]))

    # Clear all slides and re-add them via deep-copy.
    xml_slides = final_prs.slides._sldIdLst  # noqa: SLF001
    for _ in list(xml_slides):
        xml_slides.remove(xml_slides[0])

    for path in template_paths:
        src_prs = Presentation(str(path))
        for slide in src_prs.slides:
            _copy_slide(slide, final_prs)

    tmp_path = SKELETONS_DIR / f"{uuid.uuid4()}.pptx"
    final_prs.save(str(tmp_path))
    tmp_path.rename(final_path)

    logger.info(f"New skeleton saved: {final_path}")
    return str(final_path)
