from pathlib import Path
from pptx import Presentation
import hashlib
import uuid

TEMPLATE_DIR = Path("storage/templates")
GENERATED_DIR = Path("generated")
GENERATED_DIR.mkdir(exist_ok=True)


def merge_pptx(slide_names: list[str]) -> tuple[str, str]:
    """
    Merge template pptx files based on slide order.
    Returns (file_path, binary_hash)
    """

    first_template = TEMPLATE_DIR / f"{slide_names[0]}.pptx"

    if not first_template.exists():
        raise FileNotFoundError(f"{first_template} not found.")

    # İlk dosyayı base alıyoruz
    final_prs = Presentation(str(first_template))

    for slide_name in slide_names[1:]:
        template_path = TEMPLATE_DIR / f"{slide_name}.pptx"

        if not template_path.exists():
            raise FileNotFoundError(f"{template_path} not found.")

        temp_prs = Presentation(str(template_path))

        for slide in temp_prs.slides:
            final_prs.slides.add_slide(slide.slide_layout)

    temp_output = GENERATED_DIR / f"{uuid.uuid4()}.pptx"
    final_prs.save(temp_output)

    # Binary hash
    with open(temp_output, "rb") as f:
        file_hash = hashlib.sha256(f.read()).hexdigest()

    final_path = GENERATED_DIR / f"{file_hash}.pptx"
    temp_output.rename(final_path)

    return str(final_path), file_hash
